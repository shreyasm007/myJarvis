"""
Document ingestion script.

Reads documents from the data/documents directory,
chunks them, generates embeddings, and uploads to Qdrant.

Supports: TXT, MD, PDF, DOCX, DOC, ODT, RTF, HTML, CSV, JSON, XML, PPT, PPTX

Usage:
    python -m scripts.ingest
"""

import csv
import json
import os
import sys
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import List, Tuple
from uuid import uuid4

# Add project root to path
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

from dotenv import load_dotenv

# Load environment variables
load_dotenv(project_root / ".env")

# Configure proxy BEFORE importing API clients
from backend.core.proxy_config import configure_proxy
configure_proxy()

from backend.core.exceptions import DocumentIngestionError
from backend.core.logging_config import get_logger, setup_logging
from backend.rag.embeddings import get_embeddings_client
from backend.rag.vectorstore import get_vectorstore_client

# Setup logging
setup_logging("INFO")
logger = get_logger(__name__)

# Document directory
DOCUMENTS_DIR = project_root / "backend" / "data" / "documents"

# Supported file extensions
SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".markdown",  # Text
    ".pdf",  # PDF
    ".docx", ".doc", ".odt", ".rtf",  # Word processors
    ".html", ".htm",  # HTML
    ".csv", ".json", ".xml",  # Structured data
    ".ppt", ".pptx",  # PowerPoint
}

# Chunking settings
CHUNK_SIZE = 500  # characters
CHUNK_OVERLAP = 50  # characters


def read_document(file_path: Path) -> str:
    """
    Read a document file with support for multiple formats.
    
    Args:
        file_path: Path to the document
        
    Returns:
        Document content as string
    """
    try:
        ext = file_path.suffix.lower()
        
        # Text files
        if ext in {".txt", ".md", ".markdown"}:
            return read_text_file(file_path)
        
        # PDF
        elif ext == ".pdf":
            return read_pdf(file_path)
        
        # Word documents
        elif ext == ".docx":
            return read_docx(file_path)
        
        # ODT (OpenOffice)
        elif ext == ".odt":
            return read_odt(file_path)
        
        # RTF
        elif ext == ".rtf":
            return read_rtf(file_path)
        
        # HTML
        elif ext in {".html", ".htm"}:
            return read_html(file_path)
        
        # CSV
        elif ext == ".csv":
            return read_csv(file_path)
        
        # JSON
        elif ext == ".json":
            return read_json(file_path)
        
        # XML
        elif ext == ".xml":
            return read_xml(file_path)
        
        # PowerPoint
        elif ext in {".ppt", ".pptx"}:
            return read_pptx(file_path)
        
        else:
            logger.warning(f"Unsupported file type: {ext}")
            return ""
            
    except Exception as e:
        logger.error(f"Failed to read {file_path}: {str(e)}")
        raise DocumentIngestionError(
            message=f"Failed to read document: {file_path.name}",
            details={"error": str(e)},
        )


def read_text_file(file_path: Path) -> str:
    """Read plain text, markdown files."""
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def read_pdf(file_path: Path) -> str:
    """Read PDF files."""
    try:
        from pypdf import PdfReader
        
        reader = PdfReader(file_path)
        text = []
        for page in reader.pages:
            text.append(page.extract_text())
        return "\n\n".join(text)
    except ImportError:
        logger.error("pypdf not installed. Install with: pip install pypdf")
        raise


def read_docx(file_path: Path) -> str:
    """Read DOCX files."""
    try:
        from docx import Document
        
        doc = Document(file_path)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        return "\n\n".join(text)
    except ImportError:
        logger.error("python-docx not installed. Install with: pip install python-docx")
        raise


def read_odt(file_path: Path) -> str:
    """Read ODT files."""
    try:
        from odf import text as odf_text
        from odf.opendocument import load
        
        doc = load(file_path)
        paragraphs = doc.getElementsByType(odf_text.P)
        text = [str(p) for p in paragraphs if str(p).strip()]
        return "\n\n".join(text)
    except ImportError:
        logger.error("odfpy not installed. Install with: pip install odfpy")
        raise


def read_rtf(file_path: Path) -> str:
    """Read RTF files."""
    try:
        from striprtf.striprtf import rtf_to_text
        
        with open(file_path, "r", encoding="utf-8") as f:
            rtf_content = f.read()
        return rtf_to_text(rtf_content)
    except ImportError:
        logger.error("striprtf not installed. Install with: pip install striprtf")
        raise


def read_html(file_path: Path) -> str:
    """Read HTML files."""
    try:
        from bs4 import BeautifulSoup
        
        with open(file_path, "r", encoding="utf-8") as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, "lxml")
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        return soup.get_text(separator="\n\n", strip=True)
    except ImportError:
        logger.error("beautifulsoup4 not installed. Install with: pip install beautifulsoup4 lxml")
        raise


def read_csv(file_path: Path) -> str:
    """Read CSV files and convert to text."""
    with open(file_path, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        rows = []
        for row in reader:
            row_text = ", ".join([f"{k}: {v}" for k, v in row.items()])
            rows.append(row_text)
        return "\n\n".join(rows)


def read_json(file_path: Path) -> str:
    """Read JSON files and convert to text."""
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    
    def json_to_text(obj, prefix=""):
        """Recursively convert JSON to readable text."""
        if isinstance(obj, dict):
            lines = []
            for k, v in obj.items():
                if isinstance(v, (dict, list)):
                    lines.append(f"{prefix}{k}:")
                    lines.append(json_to_text(v, prefix + "  "))
                else:
                    lines.append(f"{prefix}{k}: {v}")
            return "\n".join(lines)
        elif isinstance(obj, list):
            lines = []
            for i, item in enumerate(obj):
                lines.append(f"{prefix}Item {i+1}:")
                lines.append(json_to_text(item, prefix + "  "))
            return "\n".join(lines)
        else:
            return f"{prefix}{obj}"
    
    return json_to_text(data)


def read_xml(file_path: Path) -> str:
    """Read XML files and convert to text."""
    tree = ET.parse(file_path)
    root = tree.getroot()
    
    def xml_to_text(element, prefix=""):
        """Recursively convert XML to readable text."""
        lines = []
        if element.text and element.text.strip():
            lines.append(f"{prefix}{element.tag}: {element.text.strip()}")
        else:
            lines.append(f"{prefix}{element.tag}")
        
        for child in element:
            lines.append(xml_to_text(child, prefix + "  "))
        
        return "\n".join(lines)
    
    return xml_to_text(root)


def read_pptx(file_path: Path) -> str:
    """Read PowerPoint files."""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"Slide {slide_num}:"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            text.append("\n".join(slide_text))
        
        return "\n\n".join(text)
    except ImportError:
        logger.error("python-pptx not installed. Install with: pip install python-pptx")
        raise


def chunk_text(
    text: str,
    chunk_size: int = CHUNK_SIZE,
    chunk_overlap: int = CHUNK_OVERLAP,
) -> List[str]:
    """
    Split text into overlapping chunks.
    
    Args:
        text: Text to chunk
        chunk_size: Maximum chunk size in characters
        chunk_overlap: Overlap between chunks
        
    Returns:
        List of text chunks
    """
    if not text.strip():
        return []
    
    chunks = []
    start = 0
    text_length = len(text)
    
    while start < text_length:
        end = start + chunk_size
        
        # Try to break at sentence or paragraph boundary
        if end < text_length:
            # Look for paragraph break
            para_break = text.rfind("\n\n", start, end)
            if para_break > start + chunk_size // 2:
                end = para_break + 2
            else:
                # Look for sentence break
                for sep in [". ", "! ", "? ", "\n"]:
                    sent_break = text.rfind(sep, start, end)
                    if sent_break > start + chunk_size // 2:
                        end = sent_break + len(sep)
                        break
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        # Move start with overlap
        start = end - chunk_overlap if end < text_length else text_length
    
    return chunks


def discover_documents() -> List[Tuple[Path, str]]:
    """
    Discover all documents in the documents directory.
    
    Returns:
        List of (file_path, source_name) tuples
    """
    documents = []
    
    if not DOCUMENTS_DIR.exists():
        logger.warning(f"Documents directory does not exist: {DOCUMENTS_DIR}")
        return documents
    
    for file_path in DOCUMENTS_DIR.rglob("*"):
        if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
            source_name = file_path.relative_to(DOCUMENTS_DIR).as_posix()
            documents.append((file_path, source_name))
            logger.info(f"Discovered document: {source_name}")
    
    return documents


def ingest_documents() -> int:
    """
    Main ingestion function.
    
    Reads all documents, chunks them, generates embeddings,
    and uploads to Qdrant.
    
    Returns:
        Number of chunks ingested
    """
    logger.info("Starting document ingestion")
    
    # Discover documents
    documents = discover_documents()
    
    if not documents:
        logger.warning("No documents found to ingest")
        print("\n⚠️  No documents found!")
        print(f"   Place your documents in: {DOCUMENTS_DIR}")
        print(f"   Supported formats: {', '.join(SUPPORTED_EXTENSIONS)}")
        return 0
    
    logger.info(f"Found {len(documents)} documents to ingest")
    
    # Initialize clients
    embeddings_client = get_embeddings_client()
    vectorstore_client = get_vectorstore_client()
    
    # Delete existing collection if it exists
    print("\n🗑️  Deleting existing collection...")
    logger.info("Deleting existing collection to start fresh")
    try:
        vectorstore_client.delete_collection()
        print("✅ Existing data deleted successfully")
        logger.info("Successfully deleted existing collection")
    except Exception as e:
        # Collection might not exist, which is fine
        print("ℹ️  No existing collection to delete (starting fresh)")
        logger.info(f"No existing collection found or deletion not needed: {str(e)}")
    
    # Get embedding dimension
    embedding_dim = embeddings_client.get_embedding_dimension()
    logger.info(f"Embedding dimension: {embedding_dim}")
    
    # Create fresh collection
    print("\n📦 Creating fresh collection...")
    vectorstore_client.ensure_collection_exists(embedding_dim)
    logger.info("Created new collection")
    
    # Process each document
    all_chunks = []
    all_metadata = []
    
    for file_path, source_name in documents:
        logger.info(f"Processing: {source_name}")
        
        # Read document
        content = read_document(file_path)
        
        # Chunk document
        chunks = chunk_text(content)
        logger.info(f"  Created {len(chunks)} chunks")
        
        # Add chunks with metadata
        for i, chunk in enumerate(chunks):
            all_chunks.append(chunk)
            all_metadata.append({
                "source": source_name,
                "chunk_index": i,
                "total_chunks": len(chunks),
            })
    
    if not all_chunks:
        logger.warning("No chunks created from documents")
        return 0
    
    logger.info(f"Total chunks to embed: {len(all_chunks)}")
    
    # Generate embeddings (in batches to avoid rate limits)
    batch_size = 50
    all_embeddings = []
    
    for i in range(0, len(all_chunks), batch_size):
        batch = all_chunks[i:i + batch_size]
        logger.info(f"Embedding batch {i // batch_size + 1}/{(len(all_chunks) - 1) // batch_size + 1}")
        
        embeddings = embeddings_client.embed_texts(batch, input_type="document")
        all_embeddings.extend(embeddings)
    
    # Generate IDs
    ids = [str(uuid4()) for _ in range(len(all_chunks))]
    
    # Upload to Qdrant
    logger.info("Uploading to Qdrant...")
    count = vectorstore_client.upsert_documents(
        embeddings=all_embeddings,
        documents=all_chunks,
        metadata_list=all_metadata,
        ids=ids,
    )
    
    logger.info(f"Successfully ingested {count} chunks")
    
    # Print summary
    print("\n✅ Ingestion complete!")
    print(f"   Documents processed: {len(documents)}")
    print(f"   Chunks created: {count}")
    print(f"   Collection: {vectorstore_client.collection_name}")
    
    return count


def main():
    """Main entry point."""
    print("\n📚 Document Ingestion Script")
    print("=" * 40)
    
    try:
        count = ingest_documents()
        if count > 0:
            print("\n🎉 Your documents are now ready for RAG queries!")
        sys.exit(0)
    except Exception as e:
        logger.error(f"Ingestion failed: {str(e)}")
        print(f"\n❌ Ingestion failed: {str(e)}")
        sys.exit(1)


if __name__ == "__main__":
    main()
